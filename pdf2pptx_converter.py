from pdf2image import convert_from_path
from pptx import Presentation
import sys
from io import BytesIO

def convert_pdf_to_pptx(pdf_path, pptx_path):
    try:
        images = convert_from_path(pdf_path)
        presentation = Presentation()
        for image in images:
            # Convert image data to stream
            image_stream = BytesIO()
            image.save(image_stream, format='PNG')
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.add_picture(image_stream, 0, 0)

        presentation.save(pptx_path)
        print(f"Conversion successful. PPTX file saved at {pptx_path}")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python convert_pdf_to_pptx.py <input_pdf_path> <output_pptx_path>")
    else:
        input_pdf_path = sys.argv[1]
        output_pptx_path = sys.argv[2]
        convert_pdf_to_pptx(input_pdf_path, output_pptx_path)
